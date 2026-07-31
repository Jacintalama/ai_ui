"""Markdown subset -> document blocks -> docx/pdf bytes. Pure, no I/O."""
import io

from doc_builder import parse_blocks, split_bold


def test_parse_headings_and_paragraphs():
    md = "# Title\n\nHello world.\nSame paragraph.\n\n## Section\nBody."
    b = parse_blocks(md)
    assert b[0] == {"t": "h", "level": 1, "text": "Title"}
    assert b[1] == {"t": "p", "text": "Hello world. Same paragraph."}
    assert b[2] == {"t": "h", "level": 2, "text": "Section"}
    assert b[3] == {"t": "p", "text": "Body."}


def test_parse_lists():
    md = "- one\n- two\n\n1. first\n2) second"
    b = parse_blocks(md)
    assert b[0] == {"t": "ul", "items": ["one", "two"]}
    assert b[1] == {"t": "ol", "items": ["first", "second"]}


def test_parse_code_fence_and_table():
    md = "```\nx = 1\ny = 2\n```\n\n| A | B |\n|---|---|\n| 1 | 2 |"
    b = parse_blocks(md)
    assert b[0] == {"t": "code", "text": "x = 1\ny = 2"}
    assert b[1] == {"t": "table", "rows": [["A", "B"], ["1", "2"]]}


def test_parse_never_raises_on_junk():
    assert parse_blocks("") == []
    assert parse_blocks("###### deep\n***\n> quote")  # no exception, some blocks


def test_split_bold():
    assert split_bold("a **b** c") == [("a ", False), ("b", True), (" c", False)]
    assert split_bold("plain") == [("plain", False)]
    assert split_bold("**all**") == [("all", True)]


def test_docx_roundtrip():
    from docx import Document

    from doc_builder import blocks_to_docx
    blocks = parse_blocks(
        "# Report\n\nHello **bold** world.\n\n- a\n- b\n\n| X | Y |\n|--|--|\n| 1 | 2 |")
    data = blocks_to_docx("My Report", blocks)
    assert isinstance(data, bytes) and len(data) > 2000
    doc = Document(io.BytesIO(data))
    texts = [p.text for p in doc.paragraphs]
    assert "My Report" in texts          # title heading
    assert "Report" in texts             # h1 from markdown
    assert any("Hello" in t and "world." in t for t in texts)
    bold_runs = [r.text for p in doc.paragraphs for r in p.runs if r.bold]
    assert "bold" in bold_runs
    assert "a" in texts and "b" in texts
    assert doc.tables[0].cell(0, 0).text == "X"
    assert doc.tables[0].cell(1, 1).text == "2"


def test_pdf_bytes():
    from doc_builder import blocks_to_pdf
    blocks = parse_blocks("# Report\n\nHello **bold** world.\n\n- a\n- b")
    data = blocks_to_pdf("My Report", blocks)
    assert data[:5] == b"%PDF-"
    assert len(data) > 1500
    assert b"/Page" in data


def test_pdf_escapes_markup():
    from doc_builder import blocks_to_pdf
    # < and & in user text must not break reportlab's mini-HTML parser
    data = blocks_to_pdf("t", parse_blocks("a < b & c **d**"))
    assert data[:5] == b"%PDF-"


def test_pptx_roundtrip():
    from pptx import Presentation

    from doc_builder import blocks_to_pptx
    blocks = parse_blocks(
        "# Intro\n\nHello **bold** world.\n\n- a\n- b\n\n"
        "# Data\n\n| X | Y |\n|--|--|\n| 1 | 2 |")
    data = blocks_to_pptx("My Deck", blocks)
    assert isinstance(data, bytes) and data[:2] == b"PK"   # pptx is a zip
    prs = Presentation(io.BytesIO(data))
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title]
    assert "My Deck" in titles           # title slide
    assert "Intro" in titles             # h1 -> slide
    assert "Data" in titles
    all_text = "\n".join(sh.text_frame.text for s in prs.slides
                         for sh in s.shapes if sh.has_text_frame)
    assert "Hello" in all_text and "a" in all_text and "b" in all_text
    cells = [c.text for s in prs.slides for sh in s.shapes if sh.has_table
             for row in sh.table.rows for c in row.cells]
    assert "X" in cells and "2" in cells


def test_pptx_without_headings_still_builds():
    from doc_builder import blocks_to_pptx
    data = blocks_to_pptx("t", parse_blocks("just a paragraph, no headings"))
    assert data[:2] == b"PK"

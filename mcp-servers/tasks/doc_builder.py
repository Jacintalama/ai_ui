"""Markdown subset -> blocks -> .docx / .pdf bytes.

Pure functions only; the routes layer does I/O. The subset: # ## ###
headings, paragraphs, -/* bullets, 1./1) numbered lists, **bold**, fenced
code, simple pipe tables. Anything else degrades to a paragraph; parsing
must never raise (spec: generation never errors on unknown markdown).
"""
import io
import re

_H_RE = re.compile(r"^(#{1,3})\s+(.*)$")
_UL_RE = re.compile(r"^[-*]\s+(.*)$")
_OL_RE = re.compile(r"^\d+[.)]\s+(.*)$")
_SEP_RE = re.compile(r"^\|?[\s:|-]+\|?$")


def parse_blocks(md: str) -> list:
    blocks, para, code = [], [], []
    in_code = False

    def flush_para():
        if para:
            blocks.append({"t": "p", "text": " ".join(para)})
            para.clear()

    for raw in (md or "").splitlines():
        line = raw.rstrip()
        if line.strip().startswith("```"):
            if in_code:
                blocks.append({"t": "code", "text": "\n".join(code)})
                code.clear()
            else:
                flush_para()
            in_code = not in_code
            continue
        if in_code:
            code.append(raw)
            continue
        if not line.strip():
            flush_para()
            continue
        m = _H_RE.match(line)
        if m:
            flush_para()
            blocks.append({"t": "h", "level": len(m.group(1)),
                           "text": m.group(2).strip()})
            continue
        m = _UL_RE.match(line.strip())
        if m:
            flush_para()
            if blocks and blocks[-1]["t"] == "ul":
                blocks[-1]["items"].append(m.group(1).strip())
            else:
                blocks.append({"t": "ul", "items": [m.group(1).strip()]})
            continue
        m = _OL_RE.match(line.strip())
        if m:
            flush_para()
            if blocks and blocks[-1]["t"] == "ol":
                blocks[-1]["items"].append(m.group(1).strip())
            else:
                blocks.append({"t": "ol", "items": [m.group(1).strip()]})
            continue
        if line.strip().startswith("|") and "|" in line.strip()[1:]:
            if _SEP_RE.match(line.strip()):
                continue                      # |---|---| separator row
            cells = [c.strip() for c in line.strip().strip("|").split("|")]
            flush_para()
            if blocks and blocks[-1]["t"] == "table":
                blocks[-1]["rows"].append(cells)
            else:
                blocks.append({"t": "table", "rows": [cells]})
            continue
        para.append(line.strip())
    if in_code and code:                      # unclosed fence degrades to code
        blocks.append({"t": "code", "text": "\n".join(code)})
    flush_para()
    return blocks


def split_bold(text: str) -> list:
    """'a **b** c' -> [('a ', False), ('b', True), (' c', False)]."""
    out = []
    for i, seg in enumerate((text or "").split("**")):
        if seg:
            out.append((seg, i % 2 == 1))
    return out or [("", False)]


def blocks_to_docx(title: str, blocks: list) -> bytes:
    from docx import Document

    doc = Document()
    if (title or "").strip():
        doc.add_heading(title.strip(), level=0)
    for b in blocks:
        if b["t"] == "h":
            doc.add_heading(b["text"], level=b["level"])
        elif b["t"] == "p":
            p = doc.add_paragraph()
            for seg, bold in split_bold(b["text"]):
                p.add_run(seg).bold = bold
        elif b["t"] in ("ul", "ol"):
            style = "List Bullet" if b["t"] == "ul" else "List Number"
            for it in b["items"]:
                p = doc.add_paragraph(style=style)
                for seg, bold in split_bold(it):
                    p.add_run(seg).bold = bold
        elif b["t"] == "code":
            p = doc.add_paragraph()
            run = p.add_run(b["text"])
            run.font.name = "Courier New"
        elif b["t"] == "table" and b["rows"]:
            cols = max(len(r) for r in b["rows"])
            tbl = doc.add_table(rows=len(b["rows"]), cols=cols)
            tbl.style = "Table Grid"
            for ri, row in enumerate(b["rows"]):
                for ci, cell in enumerate(row):
                    tbl.cell(ri, ci).text = cell
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def blocks_to_pptx(title: str, blocks: list) -> bytes:
    """Slides from blocks: every h1/h2 starts a new slide; p/lists/code become
    bullets on the current slide; a table gets its own slide so it never
    fights the text frame for space."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    if (title or "").strip():
        s = prs.slides.add_slide(prs.slide_layouts[0])   # title layout
        s.shapes.title.text = title.strip()

    cur = None          # current content slide's body text frame

    def new_slide(heading: str):
        nonlocal cur
        s = prs.slides.add_slide(prs.slide_layouts[1])   # title + content
        s.shapes.title.text = heading
        body = s.placeholders[1].text_frame
        body.clear()
        cur = body
        return body

    def bullet(text: str, mono: bool = False):
        body = cur if cur is not None else new_slide("")
        p = body.paragraphs[0] if (len(body.paragraphs) == 1
                                   and not body.paragraphs[0].runs) \
            else body.add_paragraph()
        for seg, bold in split_bold(text):
            r = p.add_run()
            r.text = seg
            r.font.bold = bold
            if mono:
                r.font.name = "Courier New"
                r.font.size = Pt(14)

    for b in blocks:
        if b["t"] == "h" and b["level"] <= 2:
            new_slide(b["text"])
        elif b["t"] == "h":
            bullet(f"**{b['text']}**")
        elif b["t"] == "p":
            bullet(b["text"])
        elif b["t"] in ("ul", "ol"):
            for i, it in enumerate(b["items"], 1):
                bullet(it if b["t"] == "ul" else f"{i}. {it}")
        elif b["t"] == "code":
            for line in b["text"].splitlines():
                bullet(line, mono=True)
        elif b["t"] == "table" and b["rows"]:
            s = prs.slides.add_slide(prs.slide_layouts[5])   # title only
            prev = prs.slides[-2] if len(prs.slides) > 1 else None
            s.shapes.title.text = (prev.shapes.title.text
                                   if prev is not None and prev.shapes.title
                                   else "Table")
            rows, cols = len(b["rows"]), max(len(r) for r in b["rows"])
            shp = s.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8),
                                     Inches(9), Inches(0.4 * rows))
            for ri, row in enumerate(b["rows"]):
                for ci, cell in enumerate(row):
                    shp.table.cell(ri, ci).text = cell
    if len(prs.slides) == 0:                 # no title, no headings: 1 slide
        new_slide("")
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pdf_rich(text: str) -> str:
    """Escape XML, then map **bold** to <b> for reportlab paragraphs."""
    from xml.sax.saxutils import escape
    parts = []
    for seg, bold in split_bold(text):
        seg = escape(seg)
        parts.append(f"<b>{seg}</b>" if bold else seg)
    return "".join(parts)


def blocks_to_pdf(title: str, blocks: list) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (ListFlowable, ListItem, Paragraph,
                                    Preformatted, SimpleDocTemplate, Spacer,
                                    Table, TableStyle)

    styles = getSampleStyleSheet()
    hmap = {1: "Heading1", 2: "Heading2", 3: "Heading3"}
    story = []
    if (title or "").strip():
        story += [Paragraph(_pdf_rich(title.strip()), styles["Title"]),
                  Spacer(1, 6 * mm)]
    for b in blocks:
        if b["t"] == "h":
            story.append(Paragraph(_pdf_rich(b["text"]), styles[hmap[b["level"]]]))
        elif b["t"] == "p":
            story.append(Paragraph(_pdf_rich(b["text"]), styles["BodyText"]))
        elif b["t"] in ("ul", "ol"):
            bt = "bullet" if b["t"] == "ul" else "1"
            story.append(ListFlowable(
                [ListItem(Paragraph(_pdf_rich(i), styles["BodyText"]))
                 for i in b["items"]],
                bulletType=bt))
        elif b["t"] == "code":
            story.append(Preformatted(b["text"], styles["Code"]))
        elif b["t"] == "table" and b["rows"]:
            t = Table(b["rows"], hAlign="LEFT")
            t.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ]))
            story.append(t)
        story.append(Spacer(1, 2 * mm))
    buf = io.BytesIO()
    SimpleDocTemplate(buf, pagesize=A4, title=title or "Document").build(story)
    return buf.getvalue()
